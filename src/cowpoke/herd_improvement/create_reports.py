from enum import Enum
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

from xlsxwriter import Workbook


def create_excel_spreadsheet(cow_dict: dict, xlsx_file_path: Path, csv_file_path: Path):

    workbook = Workbook(filename=xlsx_file_path)

    #################################################################################################
    # Data Summary copied from CSV file that has already been constructed.
    data_summary = workbook.add_worksheet("Data Summary")

    data_group_format = workbook.add_format()
    data_group_format.set_bg_color("#767171")
    data_group_format.set_font_color("white")

    data_header_format = workbook.add_format()
    data_header_format.set_align("center")
    data_header_format.set_bold()
    data_header_format.set_bg_color("#D0CECE")

    stats_group_format = workbook.add_format()
    stats_group_format.set_bg_color("#BF9000")
    stats_group_format.set_font_color("white")

    stats_header_format = workbook.add_format()
    stats_header_format.set_align("center")
    stats_header_format.set_bold()
    stats_header_format.set_bg_color("#FFE699")

    for col in range(21):
        data_summary.write(0, col, "", data_group_format)
    data_summary.write(0, 3, "Data from input CSV file (Lactation Data)", data_group_format)

    for col in range(21, 32):
        data_summary.write(0, col, "", stats_group_format)
    data_summary.write(0, 23, "Calculated based on input data", stats_group_format)

    row = 1

    with open(csv_file_path) as f:
        header = f.readline()
        header_fields = [x.strip() for x in header.split(",")]
        for col, field in enumerate(header_fields):
            if col < 21:
                cell_format = data_header_format
            else:
                cell_format = stats_header_format
            data_summary.write(row, col, field, cell_format)
        row += 1

        for line in f:
            values = [x.strip() for x in line.split(",")]
            for col, value in enumerate(values):
                data_summary.write(row, col, value)
            row += 1

    data_summary.set_column(first_col=0, last_col=32, width=15)

    workbook = create_report_and_report_3_8_worksheets(cow_dict=cow_dict, workbook=workbook)

    workbook.close()


def create_report_and_report_3_8_worksheets(cow_dict: dict, workbook: Workbook) -> Workbook:
    # Transform data into form useful for Report worksheets.

    cow_list = []
    for cow in cow_dict.values():
        cow_line_dict = {
            "cow_no": cow["lactation_data"]["line_number"],
            "num_lactations": cow["lactation_data"]["lact_num"],
            "milk_solids": cow["lactation_data"]["milk_solids"],
            "milk_volume": cow["lactation_data"]["milk"],
            "sire": cow["lactation_data"]["sire_name"],
            "merit_score": cow["statistics"]["merit_score"],
            "group": cow["statistics"]["group"],
        }
        cow_list.append(cow_line_dict)

    workbook = create_report_worksheet(cow_list=cow_list, workbook=workbook, worksheet_name="Report")

    cow_list_3_8 = [x.copy() for x in cow_list if x["num_lactations"] in (3, 4, 5, 6, 7, 8)]

    workbook = create_report_worksheet(cow_list=cow_list_3_8, workbook=workbook, worksheet_name="Report 3-8")

    return workbook


def create_report_worksheet(cow_list: list[dict], workbook: Workbook, worksheet_name: str) -> Workbook:

    # Calculate Ranks and Rank Comparison.

    cow_list.sort(key=lambda x: x["milk_solids"], reverse=True)
    for idx, cow in enumerate(cow_list):
        cow["milk_solids_rank"] = idx + 1

    cow_list.sort(key=lambda x: x["merit_score"], reverse=True)
    for idx, cow in enumerate(cow_list):
        cow["merit_score_rank"] = idx + 1

    for cow in cow_list:
        cow["rank_comparison"] = cow["milk_solids_rank"] - cow["merit_score_rank"]

    # Define Formats.
    bold_format = workbook.add_format({"bold": True})

    header_format = workbook.add_format()
    header_format.set_align("center")
    header_format.set_bold()
    header_format.set_bg_color("#8faadc")
    header_format.set_border()

    data_format = workbook.add_format()
    data_format.set_border()

    centred_data_format = workbook.add_format()
    centred_data_format.set_align("center")
    centred_data_format.set_border()

    databars_format = {
        "type": "data_bar", "bar_solid": True, "bar_color": "green", "bar_negative_color": "red",
        "bar_axis_position": "middle", "bar_axis_color": "black",
        # "min_type": "percent", "max_type": "percent", "min_value": "10", "max_value": "90",
    }

    field_name_to_column_name = {
        "cow_no": "Cow No",
        "num_lactations": "No of Lactations",
        "milk_solids": "305 M/S",
        "milk_volume": "305 Vol",
        "merit_score": "Merit Score",
        "milk_solids_rank": "Rank M/S",
        "merit_score_rank": "Rank Merit Score",
        "rank_comparison": "Rank Comparison",
        "group": "Group",
        "sire": "Sire",
    }
    running_max_col_widths = {
        key: len(val) for key, val in field_name_to_column_name.items()
    }
    field_order = (
        "cow_no", "num_lactations", "milk_solids", "milk_volume", "merit_score",
        "milk_solids_rank", "merit_score_rank", "rank_comparison", "group", "sire"
    )

    report = workbook.add_worksheet(worksheet_name)

    report.write(17, 0, "INDIVIDUAL ANIMAL ANALYSIS", bold_format)

    for idx, field in enumerate(field_order):
        report.write(18, idx, field_name_to_column_name[field], header_format)

    start_row = 19
    for row, cow in enumerate(cow_list):
        for col, field in enumerate(field_order):
            if field in ("rank_comparison", "sire"):
                format = data_format
            else:
                format = centred_data_format
            data_value = cow[field]
            report.write(start_row + row, col, data_value, format)

            if len(str(data_value)) > running_max_col_widths[field]:
                running_max_col_widths[field] = len(str(data_value))

    # Add databars to rank_comparison column.
    end_row = start_row + len(cow_list)
    rank_comparison_column = field_order.index("rank_comparison")
    report.conditional_format(
        first_row=start_row, first_col=rank_comparison_column, last_row=end_row, last_col=rank_comparison_column,
        options=databars_format
    )

    # Set column widths for easy viewing.
    for idx, field in enumerate(field_order):
        col_width = running_max_col_widths[field] + 4
        if col_width >= 9:  # Default column width is ~8.7 units(?)  Don't make column narrower than default.
            report.set_column(first_col=idx, last_col=idx, width=col_width)

    return workbook


class SlideLayout(Enum):
    TITLE_SLIDE = 0
    TITLE_AND_CONTENT = 1
    SECTION_HEADER = 2
    TWO_CONTENT = 3
    COMPARISON = 4
    TITLE_ONLY = 5
    BLANK = 6
    CONTENT_WITH_CAPTION = 7
    PICTURE_WITH_CAPTION = 8


def create_powerpoint(data: dict, file_path: Path):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[SlideLayout.TITLE_SLIDE.value]
    title_and_content_layout = prs.slide_layouts[SlideLayout.TITLE_AND_CONTENT.value]

    # Set size. Lengths are measured in EMU.  360,000 EMU = 1 cm.  So 1 mm = 36,000 EMU.
    # A4 page is 210 mm x 297 mm.  A4 Landscape is 297x210 mm, or 10,692,000 x 7,560,000 EMU.
    prs.slide_width = 297 * 36000  # A4 Landscape width.
    prs.slide_height = 210 * 36000  # A4 Landscape height.

    title_slide = prs.slides.add_slide(slide_layout=title_slide_layout)

    for shape in title_slide.placeholders:
        print(shape.placeholder_format.idx, shape.name)

    title = title_slide.placeholders[0]
    subtitle = title_slide.placeholders[1]
    title.top = Cm(10.5)
    subtitle.top = Cm(15)
    text_frame = title.text_frame
    para = text_frame.paragraphs[0]
    para.alignment = PP_ALIGN.RIGHT
    run = para.add_run()
    run.text = "Farm Name Here"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.name = "Assistant"
    run.font.color.rgb = RGBColor(0x12, 0x6f, 0x43)

    para = text_frame.add_paragraph()
    para.alignment = PP_ALIGN.RIGHT
    run = para.add_run()
    run.text = "Spring 2026"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.name = "Assistant"
    run.font.color.rgb = RGBColor(0x12, 0x6f, 0x43)

    para = text_frame.add_paragraph()
    para.alignment = PP_ALIGN.RIGHT
    run = para.add_run()
    run.text = "Wise Production Report"
    run.font.size = Pt(16)
    run.font.bold = False
    run.font.name = "Assistant"
    run.font.color.rgb = RGBColor(0x12, 0x6f, 0x43)

    # title.text = "Farm Name Here"
    # subtitle.text = "Wise Production Report"

    logo = title_slide.shapes.add_picture("Blockwise_Services_logo.png", left=Cm(3), top=Cm(0.5), height=Cm(10))

    slide_1 = prs.slides.add_slide(slide_layout=title_and_content_layout)
    title = slide_1.placeholders[0]
    title.text = "Wise Production Report"
    content = slide_1.placeholders[1]
    text_frame = content.text_frame
    p1 = text_frame.add_paragraph()
    p1.text = ("Utilising your milk recording data for the completed lactation,"
               " we have assessed the performance of individual animals")
    p2 = text_frame.add_paragraph()
    p2.text = ("The profitability of your herd depends on:")

    p3 = text_frame.add_paragraph()
    p3.text = "Efficiency at converting feed into milk solids."
    p3.level = 2


    prs.save(file_path)


if __name__ == "__main__":
    src_file_path = Path(__file__).parent.parent.parent
    file_path = src_file_path.joinpath("temp", "cowpoke", "herd_improvement", "blockwise_herd_report.pptx")
    create_powerpoint(data={}, file_path=file_path)
