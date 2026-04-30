from enum import Enum
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm, Pt


class SlideLayout(Enum):
    TITLE_SLIDE = 0
    TITLE_AND_CONTENT = 1
    SECTION_HEADER = 2
    TWO_CONTENT = 3
    COMPARISON = 4
    TITLE_ONLY = 5
    BLANK = 6
    CONTENT_WITH_CAPTION = 7
    PICTURE_WITH_CAPTION = 8


def create_powerpoint(data: dict, file_path: Path, image_file_paths):
    prs = Presentation()

    title_slide_layout = prs.slide_layouts[SlideLayout.TITLE_SLIDE.value]
    title_and_content_layout = prs.slide_layouts[SlideLayout.TITLE_AND_CONTENT.value]
    title_only_layout = prs.slide_layouts[SlideLayout.TITLE_ONLY.value]

    # Set size. Lengths are measured in EMU.  360,000 EMU = 1 cm.  So 1 mm = 36,000 EMU.
    # A4 page is 210 mm x 297 mm.  A4 Landscape is 297x210 mm, or 10,692,000 x 7,560,000 EMU.
    prs.slide_width = 297 * 36000  # A4 Landscape width.
    prs.slide_height = 210 * 36000  # A4 Landscape height.

    BLOCKWISE_GREEN = RGBColor(0x12, 0x6f, 0x43)

    ###########################################################################

    title_slide = prs.slides.add_slide(slide_layout=title_slide_layout)

    logo = title_slide.shapes.add_picture("Blockwise_Services_logo.png", left=Cm(3), top=Cm(0.5), height=Cm(10))

    title = title_slide.placeholders[0]
    subtitle = title_slide.placeholders[1]
    title.top = Cm(10.5)
    subtitle.top = Cm(15)

    text_frame = title.text_frame
    para = text_frame.paragraphs[0]
    para.alignment = PP_ALIGN.RIGHT
    run = para.add_run()
    run.text = "Farm Name Here"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.name = "Assistant"
    run.font.color.rgb = BLOCKWISE_GREEN

    para = text_frame.add_paragraph()
    para.alignment = PP_ALIGN.RIGHT
    run = para.add_run()
    run.text = "Spring 2026"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.name = "Assistant"
    run.font.color.rgb = BLOCKWISE_GREEN

    para = text_frame.add_paragraph()
    para.alignment = PP_ALIGN.RIGHT
    run = para.add_run()
    run.text = "Wise Production Report"
    run.font.size = Pt(16)
    run.font.bold = False
    run.font.name = "Assistant"
    run.font.color.rgb = BLOCKWISE_GREEN

    ##########################################################################################

    slide_1 = prs.slides.add_slide(slide_layout=title_and_content_layout)

    title = slide_1.placeholders[0]
    run = title.text_frame.paragraphs[0].add_run()
    run.text = "Wise Production Report"
    run.font.size = Pt(36)
    run.font.color.rgb = BLOCKWISE_GREEN

    content = slide_1.placeholders[1]

    p1 = content.text_frame.paragraphs[0]
    p1.text = ("Utilising your milk recording data for the completed lactation, "
               "we have assessed the performance of individual animals.")
    p1.runs[0].font.size = Pt(15)

    p2 = content.text_frame.add_paragraph()
    p2.text = ("The profitability of your dairy herd depends on having cows that are consistently high performers "
               "over their lifetime in the following areas:")
    p2.runs[0].font.size = Pt(15)

    p2_1 = content.text_frame.add_paragraph()
    p2_1.level = 2
    label = p2_1.add_run()
    label.text = "Production:  "
    label.font.size = Pt(15)
    label.font.bold = True
    label.font.color.rgb = BLOCKWISE_GREEN
    desc = p2_1.add_run()
    desc.text = "High production of milk solids."
    desc.font.size = Pt(15)

    p2_2 = content.text_frame.add_paragraph()
    p2_2.level = 2
    label = p2_2.add_run()
    label.text = "Efficiency:  "
    label.font.size = Pt(15)
    label.font.bold = True
    label.font.color.rgb = BLOCKWISE_GREEN
    desc = p2_2.add_run()
    desc.text = ("Efficiently converts grass into milk solids, estimated by kilograms of milk solids per "
                 "kilogram of cow liveweight.")
    desc.font.size = Pt(15)

    p2_3 = content.text_frame.add_paragraph()
    p2_3.level = 2
    label = p2_3.add_run()
    label.text = "Health:  "
    label.font.size = Pt(15)
    label.font.bold = True
    label.font.color.rgb = BLOCKWISE_GREEN
    desc = p2_3.add_run()
    desc.text = "Healthy as measured by low Somatic Cell Count (SCC)."
    desc.font.size = Pt(15)

    p2_4 = content.text_frame.add_paragraph()
    p2_4.level = 2
    label = p2_4.add_run()
    label.text = "Fertility:  "
    label.font.size = Pt(15)
    label.font.bold = True
    label.font.color.rgb = BLOCKWISE_GREEN
    desc = p2_4.add_run()
    desc.text = "In calf every year."
    desc.font.size = Pt(15)

    p3 = content.text_frame.add_paragraph()
    p3.text = ("We have used actual cow liveweight data if available.  If it is not available, we use the following "
               "standard estimates for cow liveweights:")
    p3.runs[0].font.size = Pt(15)

    p3_1 = content.text_frame.add_paragraph()
    p3_1.level = 2
    label = p3_1.add_run()
    label.text = "450 kg "
    label.font.size = Pt(15)
    label.font.bold = True
    label.font.color.rgb = BLOCKWISE_GREEN
    desc = p3_1.add_run()
    desc.text = "for cows in Lactation 1."
    desc.font.size = Pt(15)

    p3_2 = content.text_frame.add_paragraph()
    p3_2.level = 2
    label = p3_2.add_run()
    label.text = "500 kg "
    label.font.size = Pt(15)
    label.font.bold = True
    label.font.color.rgb = BLOCKWISE_GREEN
    desc = p3_2.add_run()
    desc.text = "for cows in Lactation 2 or higher."
    desc.font.size = Pt(15)

    ###########################################################################################################
    # Summary Table

    slide_2 = prs.slides.add_slide(slide_layout=title_only_layout)

    title = slide_2.placeholders[0]
    run = title.text_frame.paragraphs[0].add_run()
    run.text = "Lactation"
    run.font.size = Pt(36)
    run.font.color.rgb = BLOCKWISE_GREEN

    shape_with_table = slide_2.shapes.add_table(rows=10, cols=11, left=Cm(2), top=Cm(4), width=Cm(26), height=Cm(8))
    table = shape_with_table.table

    top_cell = table.cell(0, 0)
    top_cell.merge(table.cell(0, 10))
    top_cell.text = "Milking Herd 2025 Lactation"
    top_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    field_name_to_column_name = {
        "lactation": "Lactation",
        "num_cows": "No of Cows",
        "pct_of_herd": "% of Herd",
        "milk_volume": "Vol of Milk kg",
        "avg_milk_volume": "Avg Vol of Milk kg",
        "milk_solids": "Milk Solids kg",
        "avg_milk_solids": "Avg Milk Solids",
        "days_in_milk": "Days in Milk",
        "avg_weight": "Avg Weight",
        "avg_fat_pct": "Avg Fat %",
        "avg_protein_pct": "Avg Protein %",
    }
    running_max_col_widths = {
        key: calc_text_width(text=val, font_size=10) for key, val in field_name_to_column_name.items()
    }
    field_order = ("lactation", "num_cows", "pct_of_herd", "milk_volume", "avg_milk_volume", "milk_solids",
                   "avg_milk_solids", "days_in_milk", "avg_weight", "avg_fat_pct", "avg_protein_pct")

    for idx, field in enumerate(field_order):
        table.cell(1, idx).text = field_name_to_column_name[field]
        table.cell(1, idx).text_frame.paragraphs[0].font.size = Pt(10)
        table.cell(1, idx).text_frame.paragraphs[0].font.bold = True
        table.cell(1, idx).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # table.cell(1, idx).text_frame.paragraphs[0].font.fill.back_color = BLOCKWISE_GREEN
        table.cell(7, idx).text = field_name_to_column_name[field]
        table.cell(7, idx).text_frame.paragraphs[0].font.size = Pt(10)
        table.cell(7, idx).text_frame.paragraphs[0].font.bold = True
        table.cell(7, idx).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    group_order = ("one_and_two", "three_to_eight", "nine_plus", "total")
    for row_idx, group in enumerate(group_order):
        for idx, field in enumerate(field_order):
            data_value = str(data[group][field])
            table.cell(2 + row_idx, idx).text = data_value
            table.cell(2 + row_idx, idx).text_frame.paragraphs[0].font.size = Pt(10)
            table.cell(2 + row_idx, idx).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            if group == "total":
                table.cell(2 + row_idx, idx).text_frame.paragraphs[0].font.bold = True

            data_value_width = calc_text_width(text=data_value, font_size=10)
            if data_value_width > running_max_col_widths[field]:
                running_max_col_widths[field] = data_value_width

    group_order = ("one", "two")
    for row_idx, group in enumerate(group_order):
        for idx, field in enumerate(field_order):
            data_value = str(data[group][field])
            table.cell(8 + row_idx, idx).text = data_value
            table.cell(8 + row_idx, idx).text_frame.paragraphs[0].font.size = Pt(10)
            table.cell(8 + row_idx, idx).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            data_value_width = calc_text_width(text=data_value, font_size=10)
            if data_value_width > running_max_col_widths[field]:
                running_max_col_widths[field] = len(data_value)

    # Shrink column widths if possible.
    default_column_width = table.columns[0].width

    for idx, field in enumerate(field_order):
        col_width = running_max_col_widths[field]
        if col_width < default_column_width:
            table.columns[idx].width = col_width

    ##########################################################################################
    # Production

    slide = prs.slides.add_slide(slide_layout=title_only_layout)

    chart_file_path = str(image_file_paths.milk_solids_histogram.with_suffix(".png"))
    chart = slide.shapes.add_picture(chart_file_path, left=Cm(2), top=Cm(4), width=Cm(15))

    title = slide.placeholders[0]
    run = title.text_frame.paragraphs[0].add_run()
    run.text = "Production"
    run.font.size = Pt(36)
    run.font.color.rgb = BLOCKWISE_GREEN

    content = slide.shapes.add_textbox(left=Cm(18), top=Cm(5), height=Cm(10), width=Cm(10))
    content.text_frame.word_wrap = True

    p1 = content.text_frame.paragraphs[0]
    p1.text = ("Histogram showing number of cows in each milk production band. "
               "Eg: number of cows producing 500 to 520 kg of milk solids per year.")
    p1.runs[0].font.size = Pt(15)

    p_gap1 = content.text_frame.add_paragraph()

    p2 = content.text_frame.add_paragraph()
    p2.text = ("The distribution of cow milk solid production quite closely resembles a classic Gaussian bell curve. "
               "(As is typical for biological properties.)")
    p2.runs[0].font.size = Pt(15)

    ## Production by Age

    slide = prs.slides.add_slide(slide_layout=title_only_layout)

    chart_file_path = str(image_file_paths.milk_solids_by_age_chart.with_suffix(".png"))
    chart = slide.shapes.add_picture(chart_file_path, left=Cm(2), top=Cm(4), width=Cm(15))

    title = slide.placeholders[0]
    run = title.text_frame.paragraphs[0].add_run()
    run.text = "Production By Cow Age"
    run.font.size = Pt(36)
    run.font.color.rgb = BLOCKWISE_GREEN

    content = slide.shapes.add_textbox(left=Cm(18), top=Cm(5), height=Cm(10), width=Cm(10))
    content.text_frame.word_wrap = True

    p1 = content.text_frame.paragraphs[0]
    p1.text = ("Each dot represents an individual cow, showing its age and annual production of milk solids.")
    p1.runs[0].font.size = Pt(15)

    p_gap1 = content.text_frame.add_paragraph()

    p2 = content.text_frame.add_paragraph()
    p2.text = ("The graph shows that cows reach maximum production at about Lactation 3, and maintain production "
               "through until Lactation 8 or so.")
    p2.runs[0].font.size = Pt(15)

    ##########################################################################################
    # Efficiency

    slide = prs.slides.add_slide(slide_layout=title_only_layout)

    chart_file_path = str(image_file_paths.liveweight_milk_solids_chart.with_suffix(".png"))
    chart = slide.shapes.add_picture(chart_file_path, left=Cm(2), top=Cm(4), width=Cm(15))

    title = slide.placeholders[0]
    run = title.text_frame.paragraphs[0].add_run()
    run.text = "Efficiency"
    run.font.size = Pt(36)
    run.font.color.rgb = BLOCKWISE_GREEN

    content = slide.shapes.add_textbox(left=Cm(18), top=Cm(5), height=Cm(10), width=Cm(10))
    content.text_frame.word_wrap = True

    p1 = content.text_frame.paragraphs[0]
    p1.text = ("Each dot represents an individual cow, showing its liveweight and annual production of milk solids."
               "")
    p1.runs[0].font.size = Pt(15)

    p_gap1 = content.text_frame.add_paragraph()

    p2 = content.text_frame.add_paragraph()
    p2.text = ("The orange dots are cows for which the true liveweight is not known, so are assigned standard weight "
               "estimates of 450 kg for 1st Lactation cows and 500 kg for older cows.")
    p2.runs[0].font.size = Pt(15)

    p_gap2 = content.text_frame.add_paragraph()

    p3 = content.text_frame.add_paragraph()
    p3.text = ("The liveweight of an animal correlates with the amount of feed it consumes.  Therefore, the efficiency "
               "of a cow can be estimated by looking at the production of milk solids per kilogram of cow liveweight.")
    p3.runs[0].font.size = Pt(15)

    p_gap3 = content.text_frame.add_paragraph()

    p4 = content.text_frame.add_paragraph()
    p4.text = ("The most efficient cows produce a lot of milk for their weight, so are in the top left of the graph."
               "")
    p4.runs[0].font.size = Pt(15)

    # Efficiency vs Production

    slide = prs.slides.add_slide(slide_layout=title_only_layout)

    chart_file_path = str(image_file_paths.efficiency_milk_solids_chart.with_suffix(".png"))
    chart = slide.shapes.add_picture(chart_file_path, left=Cm(2), top=Cm(4), width=Cm(15))

    title = slide.placeholders[0]
    run = title.text_frame.paragraphs[0].add_run()
    run.text = "Efficiency and Production"
    run.font.size = Pt(36)
    run.font.color.rgb = BLOCKWISE_GREEN

    content = slide.shapes.add_textbox(left=Cm(18), top=Cm(5), height=Cm(10), width=Cm(10))
    content.text_frame.word_wrap = True

    p1 = content.text_frame.paragraphs[0]
    p1.text = ("Each dot represents an individual cow, showing its efficiency and annual production of milk solids. "
               "")
    p1.runs[0].font.size = Pt(15)

    p_gap1 = content.text_frame.add_paragraph()

    p2 = content.text_frame.add_paragraph()
    p2.text = ("The efficiency is kilograms of annual milk solids per kilogram of cow liveweight. "
               "")
    p2.runs[0].font.size = Pt(15)

    p_gap2 = content.text_frame.add_paragraph()

    p3 = content.text_frame.add_paragraph()
    p3.text = ("Orange dots are cows whose liveweights are estimated. "
               "")
    p3.runs[0].font.size = Pt(15)

    p_gap3 = content.text_frame.add_paragraph()

    p4 = content.text_frame.add_paragraph()
    p4.text = ("The best cows are most efficient so are in the right of the graph.  Efficient high producers are in "
               "the top right of the graph.")
    p4.runs[0].font.size = Pt(15)

    p_gap4 = content.text_frame.add_paragraph()

    p5 = content.text_frame.add_paragraph()
    p5.text = ("A correlation between efficiency and production can be seen.  This means the variation in milk "
               "production in the herd is more due to variation in efficiency than variation in animal size.")
    p5.runs[0].font.size = Pt(15)

    ###########################################################################################################
    # Health

    slide = prs.slides.add_slide(slide_layout=title_only_layout)

    chart_file_path = str(image_file_paths.scc_histogram.with_suffix(".png"))
    chart = slide.shapes.add_picture(chart_file_path, left=Cm(2), top=Cm(4), width=Cm(15))

    title = slide.placeholders[0]
    run = title.text_frame.paragraphs[0].add_run()
    run.text = "Health"
    run.font.size = Pt(36)
    run.font.color.rgb = BLOCKWISE_GREEN

    content = slide.shapes.add_textbox(left=Cm(18), top=Cm(5), height=Cm(10), width=Cm(10))
    content.text_frame.word_wrap = True

    p1 = content.text_frame.paragraphs[0]
    p1.text = ("Histogram showing percentage of cows in each band of average somatic cell count (SCC). "
               "Eg: percentage of cows with an Average SCC from 50 to 100.")
    p1.runs[0].font.size = Pt(15)

    p_gap1 = content.text_frame.add_paragraph()

    p2 = content.text_frame.add_paragraph()
    p2.text = ("The highly-skewed distribution shows that the vast majority of cows are very healthy. "
               "There are a few outliers with extremely high Average SCCs.")
    p2.runs[0].font.size = Pt(15)

    # Health and Production

    slide = prs.slides.add_slide(slide_layout=title_only_layout)

    chart_file_path = str(image_file_paths.milk_solids_vs_scc_chart.with_suffix(".png"))
    chart = slide.shapes.add_picture(chart_file_path, left=Cm(2), top=Cm(4), width=Cm(15))

    title = slide.placeholders[0]
    run = title.text_frame.paragraphs[0].add_run()
    run.text = "Health and Production"
    run.font.size = Pt(36)
    run.font.color.rgb = BLOCKWISE_GREEN

    content = slide.shapes.add_textbox(left=Cm(18), top=Cm(5), height=Cm(10), width=Cm(10))
    content.text_frame.word_wrap = True

    p1 = content.text_frame.paragraphs[0]
    p1.text = ("Each dot represents an individual cow, showing its average somatic cell count (SCC) and annual milk "
               "solids production.")
    p1.runs[0].font.size = Pt(15)

    p_gap1 = content.text_frame.add_paragraph()

    p2 = content.text_frame.add_paragraph()
    p2.text = ("Interestingly, poor health as measured by a high average SCC does not seem to be particularly "
               "correlated with lower milk production.")
    p2.runs[0].font.size = Pt(15)

    # Health and Efficiency

    slide = prs.slides.add_slide(slide_layout=title_only_layout)

    chart_file_path = str(image_file_paths.efficiency_vs_scc_chart.with_suffix(".png"))
    chart = slide.shapes.add_picture(chart_file_path, left=Cm(2), top=Cm(4), width=Cm(15))

    title = slide.placeholders[0]
    run = title.text_frame.paragraphs[0].add_run()
    run.text = "Health and Efficiency"
    run.font.size = Pt(36)
    run.font.color.rgb = BLOCKWISE_GREEN

    content = slide.shapes.add_textbox(left=Cm(18), top=Cm(5), height=Cm(10), width=Cm(10))
    content.text_frame.word_wrap = True

    p1 = content.text_frame.paragraphs[0]
    p1.text = ("Each dot represents an individual cow, showing its average somatic cell count (SCC) and efficiency, "
               "which is kilograms of milk solids per kilogram of liveweight.")
    p1.runs[0].font.size = Pt(15)

    p_gap1 = content.text_frame.add_paragraph()

    p2 = content.text_frame.add_paragraph()
    p2.text = ("Interestingly, poor health as measured by a high average SCC does not seem to be particularly "
               "correlated with lower efficiency at converting feed into milk.")
    p2.runs[0].font.size = Pt(15)

    ##################################################################################
    # Cow Weights

    slide = prs.slides.add_slide(slide_layout=title_only_layout)

    chart_file_path = str(image_file_paths.liveweight_histogram.with_suffix(".png"))
    chart = slide.shapes.add_picture(chart_file_path, left=Cm(2), top=Cm(4), width=Cm(15))

    title = slide.placeholders[0]
    run = title.text_frame.paragraphs[0].add_run()
    run.text = "Herd Liveweights"
    run.font.size = Pt(36)
    run.font.color.rgb = BLOCKWISE_GREEN

    content = slide.shapes.add_textbox(left=Cm(18), top=Cm(5), height=Cm(10), width=Cm(10))
    content.text_frame.word_wrap = True

    p1 = content.text_frame.paragraphs[0]
    p1.text = ("Histogram showing percentage of herd in each weight band. "
               "Eg: percentage of herd weighing 500 kg to 525 kg.")
    p1.runs[0].font.size = Pt(15)

    p_gap1 = content.text_frame.add_paragraph()

    p2 = content.text_frame.add_paragraph()
    p2.text = ("The distribution of cow weights quite closely resembles a classic Gaussian bell curve. "
               "(As is typical for biological properties.)")
    p2.runs[0].font.size = Pt(15)

    ############################################################################
    # Age Distribution

    slide = prs.slides.add_slide(slide_layout=title_only_layout)

    chart_file_path = str(image_file_paths.cow_age_chart.with_suffix(".png"))
    chart = slide.shapes.add_picture(chart_file_path, left=Cm(2), top=Cm(4), width=Cm(15))

    title = slide.placeholders[0]
    run = title.text_frame.paragraphs[0].add_run()
    run.text = "Herd Age Distribution"
    run.font.size = Pt(36)
    run.font.color.rgb = BLOCKWISE_GREEN

    content = slide.shapes.add_textbox(left=Cm(18), top=Cm(5), height=Cm(10), width=Cm(10))
    content.text_frame.word_wrap = True

    p1 = content.text_frame.paragraphs[0]
    p1.text = ("Bar chart showing percentage of herd at each age (lactation).")
    p1.runs[0].font.size = Pt(15)

    p_gap1 = content.text_frame.add_paragraph()

    p2 = content.text_frame.add_paragraph()
    p2.text = ("The red diamonds joined by red lines are the industry standard age distribution.")
    p2.runs[0].font.size = Pt(15)

    ############################################################################
    # Add green triangle decorations.

    for slide in prs.slides:
        h = 13.45  # Originally 15.74
        b = 2.0  # Originally 2.34
        top = h/2 - b/2
        left = -top
        triangle =slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, left=Cm(left), top=Cm(top), height=Cm(b), width=Cm(h))
        triangle.rotation = 90
        triangle.line.color.rgb = BLOCKWISE_GREEN
        triangle.fill.solid()
        triangle.fill.fore_color.rgb = BLOCKWISE_GREEN

        h = 8.41  # Originally 7.90.  Changed to 8.41 to match shape of upper triangle
        b = 1.25
        top = 21 - h/2 - b/2  # A4 page portrait height = 21.0 cm.
        left = 29.7 - h/2 - b/2  # A4 page portrait width = 29.7 cm.
        triangle =slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, left=Cm(left), top=Cm(top), height=Cm(b), width=Cm(h))
        triangle.rotation = -90
        triangle.line.color.rgb = BLOCKWISE_GREEN
        triangle.fill.solid()
        triangle.fill.fore_color.rgb = BLOCKWISE_GREEN

    prs.save(file_path)


def calc_text_width(text: str, font_size: int) -> int:
    """Estimate width of string in EMU.  Adds padding of one character either side."""
    return int((len(text) + 2) * Pt(font_size) * 0.53)


if __name__ == "__main__":
    src_file_path = Path(__file__).parent.parent.parent
    file_path = src_file_path.joinpath("temp", "cowpoke", "herd_improvement", "blockwise_herd_report.pptx")

    from lactation_calculations import FilePaths

    group_data = {
        "lactation": "1 & 2",
        "num_cows": 300,
        "pct_of_herd": 75,
        "milk_volume": 5123456,
        "avg_milk_volume": 5000,
        "milk_solids": 12341,
        "avg_milk_solids": 234,
        "days_in_milk": 300,
        "avg_weight": 550,
        "avg_fat_pct": 4.3,
        "avg_protein_pct": 3.7,
    }
    summary_stats = {
        "one_and_two": group_data,
        "three_to_eight": group_data,
        "nine_plus": group_data,
        "total": group_data,
        "one": group_data,
        "two": group_data,
    }
    create_powerpoint(data=summary_stats, file_path=file_path, image_file_paths=FilePaths)
